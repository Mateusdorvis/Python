#CRIANDO UM SLIDE COM PYTHON
from pptx import Presentation
from pptx.util import Inches

apresentacao = Presentation()

#cria o primeiro slide
slide_titulo = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])
titulo = slide_titulo.shapes.title
subtitulo =  slide_titulo.placeholders[1]
#inserindo o texto e o subtitulo no slide 1
titulo.text = 'Slde com imagem'
subtitulo.text = 'Este slide é com imagem !'

#criando o segundo slide 
conteudo_slide =apresentacao.slides.add_slide(apresentacao.slide_layouts[1])
titulo_slide =  conteudo_slide.shapes.title

conteudo_slide.shapes.add_picture(r'C:\xampp\htdocs\progr\python\ARTE\skylite\oi.jpg', Inches(1), Inches(1), width=Inches(5)) #insere a imagem no slide


titulo_slide.text = 'A imagem abaixo é maravilhosa!'

#salvando o slide completo
apresentacao.save('slide_com_imagem.pptx')
print('Slide criado !')
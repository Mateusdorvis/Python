#CRIANDO UM SLIDE COM PYTHON
from pptx import Presentation
from pptx.util import Inches

apresentacao = Presentation()

#cria o primeiro slide
slide_titulo = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])
titulo = slide_titulo.shapes.title
subtitulo =  slide_titulo.placeholders[1]
#inserindo o texto e o subtitulo no slide 1
titulo.text = 'Slide com python'
subtitulo.text = 'Este slide eu fiz com python'

#criando o segundo slide 
conteudo_slide =apresentacao.slides.add_slide(apresentacao.slide_layouts[1])
titulo_slide =  conteudo_slide.shapes.title
subtitulo_slide = conteudo_slide.placeholders[1]


titulo_slide.text = 'A linguagem python é maravilhosa!'
subtitulo_slide.text = 'A  linguagem em si é tão incrível que fez este slide !'
#salvando o slide completo
apresentacao.save('meu_slide.pptx')
print('Slide criado !')